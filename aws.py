from pptx import Presentation

# Create a PowerPoint presentation object
presentation = Presentation()

# Add a title slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "AWS SAA-C03 Practice Questions"
subtitle.text = "Prepare for the AWS Solutions Architect Associate Exam"

# List of questions
questions = [
    {
        "question": "You are designing a web application that will be hosted on Amazon EC2 instances behind an Application Load Balancer (ALB). The application needs to store session data for users. Which AWS service is the most cost-effective and scalable solution for storing session data?",
        "options": ["A. Amazon RDS", "B. Amazon DynamoDB", "C. Amazon ElastiCache", "D. Amazon S3"]
    },
    {
        "question": "A company is migrating its on-premises application to AWS. The application requires a shared file system that multiple EC2 instances can access concurrently. Which AWS service should you recommend?",
        "options": ["A. Amazon S3", "B. Amazon EFS", "C. Amazon FSx for Windows File Server", "D. Amazon RDS"]
    },
    {
        "question": "Your company wants to use Amazon S3 to store sensitive data. Compliance requirements mandate that the data must be encrypted at rest. Which of the following options will meet this requirement? (Choose TWO.)",
        "options": ["A. Use Amazon S3 default encryption with SSE-S3.", "B. Use Amazon S3 default encryption with SSE-KMS.", "C. Use Amazon S3 default encryption with SSE-C.", "D. Enable versioning on the S3 bucket.", "E. Use Amazon S3 Transfer Acceleration."]
    },
    {
        "question": "A company is running a web application on AWS using Amazon RDS for its database. The application experiences high read traffic, which is causing performance issues. What is the best way to improve the database's read performance?",
        "options": ["A. Enable Multi-AZ deployment for the RDS instance.", "B. Create a read replica of the RDS instance.", "C. Use Amazon ElastiCache in front of the RDS instance.", "D. Increase the size of the RDS instance."]
    },
    {
        "question": "You are designing a solution that requires a highly available database. The database must automatically failover to a standby instance in another Availability Zone in case of an outage. Which AWS service should you use?",
        "options": ["A. Amazon Aurora", "B. Amazon DynamoDB", "C. Amazon RDS with Multi-AZ", "D. Amazon Redshift"]
    },{
        "question": "A company is using AWS Lambda functions to process data from an Amazon S3 bucket. The Lambda function needs to access the S3 bucket, but you want to ensure that the function has the least privilege necessary. What is the best way to grant the Lambda function access to the S3 bucket?",
        "options": ["A. Attach an IAM role with S3 permissions to the Lambda function.", "B. Use an S3 bucket policy to allow access to the Lambda function.", "C. Use an IAM user with S3 permissions for the Lambda function.", "D. Use AWS Secrets Manager to store S3 credentials for the Lambda function."]
    }
]

# Add slides for each question
for i, q in enumerate(questions, start=1):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = f"Question {i}"
    content.text = q["question"] + "\n\n" + "\n".join(q["options"])

# Save the presentation
presentation.save("AWS_SAA_C03_Practice_Questions.pptx")
print("Presentation saved as 'AWS_SAA_C03_Practice_Questions.pptx'")