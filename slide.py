from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation for detailed DSA in Java
prs_detailed = Presentation()

# Helper function to add slide with title and content
def add_detailed_slide(title, content):
    slide = prs_detailed.slides.add_slide(prs_detailed.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Title slide
slide_title = prs_detailed.slides.add_slide(prs_detailed.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]
title.text = "Data Structures and Algorithms (DSA) in Java \n Lecture Notes \t by MG"
subtitle.text = "Detailed Lecture Notes and Examples"

# Add detailed slides
add_detailed_slide("Introduction to DSA",
    "• DSA is the core of effective programming.\n• Data Structures: organize and manage data.\n• Algorithms: define logic to manipulate data.\n\n"
    "Benefits:\n- Improve performance\n- Scalability\n- Master coding interviews and real-world problem solving.")

add_detailed_slide("Arrays in Java",
    "• Fixed-size container holding elements of the same type.\n• Indexed access starting from 0.\n\n"
    "Syntax:\nint[] arr = new int[5];\n\n"
    "Common Operations:\n- Access: O(1)\n- Search: O(n)\n- Insert/Delete: O(n)\n\nUse Cases:\n- Lookup tables\n- Image processing")

add_detailed_slide("Singly Linked List",
    "• Linear data structure where elements (nodes) point to the next.\n\n"
    "Node Example:\nclass Node {\n  int data;\n  Node next;\n  Node(int d) { data = d; next = null; }\n}\n\n"
    "Operations:\n- Insert, Delete, Search\n- Traversal using loops\n\nUse Cases:\n- Dynamic memory\n- Undo/Redo features")

add_detailed_slide("Stacks and Queues",
    "• Stack: LIFO\n  - Methods: push(), pop(), peek()\n  - Example: Call stack, Undo feature\n\n"
    "• Queue: FIFO\n  - Methods: offer(), poll(), peek()\n  - Example: Job scheduling, Messaging\n\n"
    "Use Java's Stack and LinkedList classes to implement.")

add_detailed_slide("Binary Search Tree (BST)",
    "• Tree where left < root < right.\n\nNode:\nclass Node {\n  int key;\n  Node left, right;\n}\n\n"
    "Operations:\n- Insertion, Search, Deletion\n- Inorder Traversal = Sorted Order\n\nUse Cases:\n- Indexing in databases\n- Auto-complete")

add_detailed_slide("HashMap and HashSet",
    "• HashMap: Key-Value pairs.\n• HashSet: Unique values only (internally uses HashMap).\n\n"
    "Java:\nMap<String, Integer> map = new HashMap<>();\nSet<String> set = new HashSet<>();\n\n"
    "Use Cases:\n- Caching\n- Fast data access\n- Removing duplicates")

add_detailed_slide("Graphs in Java",
    "• Set of vertices connected by edges.\n• Directed/Undirected, Weighted/Unweighted.\n\n"
    "Representations:\n- Adjacency Matrix\n- Adjacency List\n\n"
    "Use Cases:\n- Navigation systems\n- Social networks\n- Recommendations")

add_detailed_slide("Sorting Algorithms",
    "• Bubble Sort: O(n²), simple but slow.\n• Merge Sort: O(n log n), stable.\n• Quick Sort: O(n log n), fast but not stable.\n\n"
    "Use Cases:\n- Preparing data for search\n- Efficient display\n- Analytics and processing")

add_detailed_slide("Searching Algorithms",
    "• Linear Search: O(n), unsorted data.\n• Binary Search: O(log n), sorted arrays only.\n\n"
    "Code Example (Binary Search):\nwhile (low <= high) {\n  int mid = (low + high) / 2;\n  if (arr[mid] == key) return mid;\n  ...\n}\n\n"
    "Use Cases:\n- Name lookups\n- ID verification")

add_detailed_slide("Conclusion",
    "• Mastering DSA improves coding, logic, and design.\n• Helps in interviews, system development, and optimization.\n\n"
    "Recommended Practice Platforms:\n- LeetCode\n- HackerRank\n- GeeksforGeeks\n\nKeep coding, keep learning!")

# Save detailed presentation
pptx_detailed_path = r"D:\works\MG-DSA\DSA_in_Java_LectureNotes_Detailed.pptx"
prs_detailed.save(pptx_detailed_path)

pptx_detailed_path
