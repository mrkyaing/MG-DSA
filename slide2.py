from pptx import Presentation
from pptx.util import Pt

# Load the existing presentation
prs = Presentation(r"D:\works\MG-DSA\DSA_in_Java_LectureNotes2.pptx")

# Function to set detailed content for a given slide index
def update_slide_content(slide_index, content):
    slide = prs.slides[slide_index]
    slide.placeholders[1].text = content

# Update slides with more detailed content
update_slide_content(1,  # Introduction to DSA
    "• Data Structures and Algorithms (DSA) form the basis of problem-solving in computer science.\n"
    "• Data Structures: Efficient storage, organization, and management of data.\n"
    "• Algorithms: Well-defined steps to solve a problem.\n\n"
    "Benefits:\n"
    "- Optimized code performance\n"
    "- Scalability\n"
    "- Essential in interviews and system design")

update_slide_content(2,  # Arrays in Java
    "Arrays are fixed-size sequential data structures used to store elements of the same type.\n\n"
    "Syntax:\n"
    "int[] arr = new int[5];\n"
    "arr[0] = 10;\n\n"
    "Operations:\n"
    "- Traversal: for/while loop\n"
    "- Searching: Linear/Binary\n"
    "- Sorting: Bubble, Merge, Quick\n\n"
    "Use Cases:\n"
    "- Static data storage\n"
    "- Lookup tables")

update_slide_content(3,  # Singly Linked List
    "A linked list is a linear data structure where each element points to the next.\n\n"
    "Node class:\n"
    "class Node {\n"
    "    int data;\n"
    "    Node next;\n"
    "    Node(int data) { this.data = data; this.next = null; }\n"
    "}\n\n"
    "LinkedList Operations:\n"
    "- add(), remove(), search(), display()\n\n"
    "Advantages:\n"
    "- Dynamic size\n"
    "- Efficient insert/delete\n\n"
    "Disadvantages:\n"
    "- No direct access (unlike arrays)")

update_slide_content(4,  # Stacks and Queues
    "Stack: Last In First Out (LIFO)\n"
    "Methods: push(), pop(), peek()\n"
    "Java: Stack<Integer> stack = new Stack<>();\n\n"
    "Queue: First In First Out (FIFO)\n"
    "Methods: offer(), poll(), peek()\n"
    "Java: Queue<Integer> queue = new LinkedList<>();\n\n"
    "Applications:\n"
    "- Stack: Undo feature, call stack\n"
    "- Queue: Thread scheduling, BFS traversal")

update_slide_content(5,  # BST
    "A binary search tree (BST) stores nodes in a sorted manner for fast lookup and modification.\n\n"
    "Rules:\n"
    "- Left child < parent\n"
    "- Right child > parent\n\n"
    "Operations:\n"
    "- Insertion, Deletion, Search\n"
    "- InOrder, PreOrder, PostOrder Traversal\n\n"
    "Time Complexity:\n"
    "- Best/Average: O(log n)\n"
    "- Worst (unbalanced): O(n)\n\n"
    "Real Use Cases:\n"
    "- Databases, Compiler Syntax Trees")

update_slide_content(6,  # HashMap and HashSet
    "HashMap: Stores key-value pairs.\n"
    "HashSet: Stores unique elements (backed by HashMap).\n\n"
    "HashMap Example:\n"
    "Map<String, Integer> map = new HashMap<>();\n"
    "map.put(\"A\", 1);\n\n"
    "HashSet Example:\n"
    "Set<String> set = new HashSet<>();\n"
    "set.add(\"apple\");\n\n"
    "Applications:\n"
    "- Caching\n"
    "- Counting frequency\n"
    "- Fast lookups")

update_slide_content(7,  # Graphs
    "Graphs are structures with nodes (vertices) and connections (edges).\n\n"
    "Representations:\n"
    "- Adjacency Matrix (2D array)\n"
    "- Adjacency List (HashMap of lists)\n\n"
    "Traversals:\n"
    "- BFS (Queue)\n"
    "- DFS (Stack/Recursion)\n\n"
    "Use Cases:\n"
    "- Pathfinding (shortest path)\n"
    "- Network routing\n"
    "- Social media connections")

update_slide_content(8,  # Sorting Algorithms
    "Sorting arranges data in a specific order (ascending/descending).\n\n"
    "Common Sorting Algorithms:\n"
    "- Bubble Sort: O(n²), simple, inefficient\n"
    "- Selection Sort: O(n²), selects min/max\n"
    "- Merge Sort: O(n log n), divide and conquer\n"
    "- Quick Sort: O(n log n) average, O(n²) worst\n\n"
    "Applications:\n"
    "- Organizing data\n"
    "- Preparing for binary search\n"
    "- Efficient record keeping")

update_slide_content(9,  # Searching Algorithms
    "Linear Search:\n"
    "- Traverse each element\n"
    "- O(n) time complexity\n\n"
    "Binary Search:\n"
    "- Requires sorted array\n"
    "- Repeatedly divides search range\n"
    "- O(log n) time complexity\n\n"
    "Applications:\n"
    "- Search engines\n"
    "- Database queries\n"
    "- Games")

update_slide_content(10,  # Conclusion
    "• DSA is essential for every developer.\n"
    "• Strengthens logical and analytical thinking.\n"
    "• Real-world relevance in systems, databases, applications.\n\n"
    "Resources to Practice:\n"
    "- LeetCode (interview-style problems)\n"
    "- HackerRank (certifications and contests)\n"
    "- GeeksforGeeks (detailed explanations)\n\n"
    "Advice:\n"
    "- Code daily\n"
    "- Solve different problem types\n"
    "- Participate in coding challenges")

# Save updated presentation
prs = Presentation(r"D:\works\MG-DSA\DSA_in_Java_LectureNotes2.pptx")
prs.save(detailed_pptx_path)

detailed_pptx_path
