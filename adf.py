from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Create a presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ETL vs ELT in Azure Data Factory"
subtitle.text = "Core Concepts, Differences, and Use Cases"

# Slide 1: ETL vs ELT Overview
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "ETL vs ELT: Main Concepts"
content = slide.placeholders[1]
content.text = (
    "ETL: Extract → Transform → Load\n"
    "ELT: Extract → Load → Transform\n"
    "ETL transforms data before loading it to target\n"
    "ELT loads raw data first and transforms within the target system"
)

# Slide 2: Comparison Table
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "ETL vs ELT: Key Differences"

# Create a table
rows, cols = 7, 3
left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(3.5)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column headers
table.cell(0, 0).text = "Feature"
table.cell(0, 1).text = "ETL"
table.cell(0, 2).text = "ELT"

# Fill rows
data = [
    ("Order", "Extract → Transform → Load", "Extract → Load → Transform"),
    ("Transform location", "Outside target system", "Inside target system"),
    ("Target system", "Traditional data warehouse", "Modern cloud warehouse"),
    ("Performance", "May be slower with big data", "Faster with large data"),
    ("Best for", "Legacy systems", "Modern platforms like Synapse"),
    ("Example", "SSIS + SQL Server", "ADF + Azure Synapse"),
]

for i, (feature, etl, elt) in enumerate(data, start=1):
    table.cell(i, 0).text = feature
    table.cell(i, 1).text = etl
    table.cell(i, 2).text = elt

# Slide 3: ETL Use Case
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "ETL Use Case Example"
content = slide.placeholders[1]
content.text = (
    "A bank extracts data from multiple systems.\n"
    "Cleans and transforms data using SSIS.\n"
    "Loads cleaned data into a traditional data warehouse.\n"
    "Used for financial reporting and auditing."
)

# Slide 4: ELT Use Case
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "ELT Use Case Example"
content = slide.placeholders[1]
content.text = (
    "An e-commerce site collects user clickstream data.\n"
    "Loads raw logs into Azure Data Lake.\n"
    "Transforms data using Azure Synapse SQL or ADF Data Flows.\n"
    "Supports real-time analytics and ML models."
)

# Slide 5: When to Use ETL or ELT in ADF
slide = prs.slides.add_slide(bullet_slide_layout)
slide.shapes.title.text = "ETL or ELT in Azure Data Factory?"
content = slide.placeholders[1]
content.text = (
    "ADF supports both ETL and ELT.\n"
    "Use ELT for cloud-native platforms (e.g., Synapse, Data Lake).\n"
    "Use ETL when transformation must happen before loading.\n"
    "Choose based on system capabilities and data volume."
)

# Save presentation
pptx_path = "D:\\works\\MG-DSA\\ETL_vs_ELT_Azure_Data_Factory.pptx"
prs.save(pptx_path)

pptx_path
