from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Create a presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Add title slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Software Development Methodologies Overview"
subtitle.text = "Waterfall, Agile, Scrum, and Jira (In Burmese)"

# Slide content (each section becomes a slide)
slides_content = [
    {
        "title": "ဆော့ဖ်ဝဲ ဖွံ့ဖြိုးရေး Methodology များ",
        "content": "Software Development Methodology ဆိုသည်မှာ ဆော့ဖ်ဝဲတစ်ခုကို စတင်ဖန်တီးသည်မှစ၍ ပြီးဆုံးလာသည်အထိ "
                   "လုပ်ငန်းစဉ်များကို စနစ်တကျ ဖန်တီးလုပ်ဆောင်ရာတွင် အသုံးပြုသော နည်းလမ်းများဖြစ်သည်။"
    },
    {
        "title": "Waterfall Methodology",
        "content": "Waterfall သည် ရိုးရှင်းပြီး အဆင့်လိုက်ဆင်းသက်မှုရှိသော နည်းလမ်းဖြစ်သည်။\n\nSteps:\n1. Requirements\n2. Design\n"
                   "3. Implementation\n4. Testing\n5. Deployment\n6. Maintenance\n\nသင့်လျော်သောအခြေအနေများ:\n- လိုအပ်ချက်များ ပြီးပြည့်စုံသော စီမံကိန်းများ\n"
                   "- ပြောင်းလဲမှုနည်းသောလုပ်ငန်းများ"
    },
    {
        "title": "Agile Methodology",
        "content": "Agile သည် တိုးတက်မှုကို Iteration အတိုင်း အဆင့်လိုက်လုပ်ဆောင်ပြီး ပြောင်းလဲမှုများအပေါ် တုန့်ပြန်နိုင်သော နည်းလမ်းဖြစ်သည်။\n\nအဓိကလက္ခဏာများ:\n"
                   "- Iteration-based (Sprints)\n- Customer Feedback\n- Team Collaboration\n- Adaptability"
    },
    {
        "title": "Scrum Framework",
        "content": "Scrum သည် Agile ၏ အစိတ်အပိုင်းဖြစ်ပြီး Sprint များအတိုင်း ဖွဲ့စည်းသည်။\n\nRoles:\n- Product Owner\n- Scrum Master\n"
                   "- Development Team\n\nEvents:\n- Sprint Planning\n- Daily Scrum\n- Sprint Review\n- Sprint Retrospective"
    },
    {
        "title": "Jira Tool",
        "content": "Jira သည် Agile methodology အတွက် အသုံးပြုနိုင်သော Project Management Tool တစ်ခုဖြစ်သည်။\n\nအသုံးများသောအရာများ:\n"
                   "- Task Assignment\n- Sprint Planning\n- Issue Tracking\n\nAgile & Scrum အသုံးပြုသည့်အဖွဲ့များအတွက် အထောက်အကူဖြစ်စေသည်။"
    },
    {
        "title": "နိဂုံးချုပ်",
        "content": "Methodology များကို လုပ်ငန်းလိုအပ်ချက်နှင့်အညီ 適切 にရွေးချယ်ရန် လိုအပ်သည်။\n\n- Waterfall = တိကျမှုလိုသောစီမံကိန်းများ\n"
                   "- Agile/Scrum = ပြောင်းလဲမှုများနှင့် ကိုက်ညီသောစီမံကိန်းများ"
    }
]

# Add content slides
for item in slides_content:
    slide = prs.slides.add_slide(content_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = item["title"]
    content.text = item["content"]

# Save the presentation
pptx_path = r"D:\works\MG-DSA\SE.pptx"

prs.save(pptx_path)

pptx_path
